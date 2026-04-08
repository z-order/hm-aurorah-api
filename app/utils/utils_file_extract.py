"""
File text extraction utilities

Extracts plain text from various file formats. Each extractor inserts
double newlines (\\n\\n) at logical boundaries (paragraphs, slides, rows, pages).
The returned text is then passed through analyze_raw_text_to_json() +
add_sentence_markers() for final sentence-level segmentation.
"""

import io
import json
import logging
import xml.etree.ElementTree as ET
import zipfile
from typing import cast

import fitz  # type: ignore[import-untyped]
import httpx
from docx import Document
from google.api_core.client_options import ClientOptions  # type: ignore[import-untyped]
from google.cloud import documentai  # type: ignore[import-untyped]
from google.oauth2 import service_account  # type: ignore[import-untyped]
from openpyxl import load_workbook
from pptx import Presentation
from striprtf.striprtf import rtf_to_text  # type: ignore[import-untyped]

from app.core.config import settings
from app.core.logger import get_logger
from app.utils.utils_http import decode_bytes

logger = get_logger(__name__, logging.INFO)


# =============================================================================
# DOCX  (python-docx)
# =============================================================================


def extract_text_from_docx(file_bytes: bytes) -> str:
    """Extract text from DOCX. Paragraphs separated by \\n\\n."""
    doc = Document(io.BytesIO(file_bytes))  # type: ignore[no-untyped-call]
    paragraphs: list[str] = []
    for para in doc.paragraphs:  # type: ignore[union-attr]
        t = str(para.text).strip()  # type: ignore[union-attr]
        if t:
            paragraphs.append(t)

    return "\n\n".join(paragraphs)


# =============================================================================
# PPTX  (python-pptx)
# =============================================================================


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """Extract text from PPTX. Slides separated by \\n\\n."""
    prs = Presentation(io.BytesIO(file_bytes))
    slide_texts: list[str] = []

    for slide in prs.slides:  # type: ignore[union-attr]
        parts: list[str] = []
        for shape in slide.shapes:  # type: ignore[union-attr]
            if shape.has_text_frame:  # type: ignore[union-attr]
                for paragraph in shape.text_frame.paragraphs:  # type: ignore[union-attr]
                    t = str(paragraph.text).strip()  # type: ignore[union-attr]
                    if t:
                        parts.append(t)
        if parts:
            slide_texts.append("\n".join(parts))

    return "\n\n".join(slide_texts)


# =============================================================================
# XLSX  (openpyxl)
# =============================================================================


def extract_text_from_xlsx(file_bytes: bytes) -> str:
    """Extract text from XLSX. Rows separated by \\n\\n."""
    wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    row_texts: list[str] = []

    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                row_texts.append("\t".join(cells))

    wb.close()
    return "\n\n".join(row_texts)


# =============================================================================
# HWPX  (ZIP + XML)
# =============================================================================

_HWPX_NS = {
    "hp": "http://www.hancom.co.kr/hwpml/2011/paragraph",
    "hp2": "urn:hancom:hwpml:2011",
}


def extract_text_from_hwpx(file_bytes: bytes) -> str:
    """
    Extract text from HWPX. Paragraphs separated by \\n\\n.

    Reads Contents/section0.xml, section1.xml, ... in order,
    extracts text from <hp:t> tags.
    """
    paragraphs: list[str] = []

    with zipfile.ZipFile(io.BytesIO(file_bytes)) as zf:
        section_files = sorted(n for n in zf.namelist() if n.startswith("Contents/section") and n.endswith(".xml"))

        for section_file in section_files:
            xml_data = zf.read(section_file)
            root = ET.fromstring(xml_data)

            for t_elem in root.iter():
                if t_elem.tag.endswith("}t") or t_elem.tag == "t":
                    text = (t_elem.text or "").strip()
                    if text:
                        paragraphs.append(text)

    return "\n\n".join(paragraphs)


# =============================================================================
# PDF  (PyMuPDF + Google OCR(Document AI) + Upstage OCR fallback)
#
# Google OCR(Document AI): https://docs.cloud.google.com/document-ai/docs/overview
# Upstage OCR: https://console.upstage.ai/docs/getting-started?api=document-parsing
# =============================================================================


def extract_text_from_pdf(file_bytes: bytes, scan_mode: str = "text") -> str:
    """
    Extract text from PDF using the chosen scan mode.

    Args:
        file_bytes: Raw PDF bytes
        scan_mode: "text" (PyMuPDF direct), "google_ocr" (Google Document AI),
                   or "upstage_ocr" (Upstage Document Digitization)
    """
    if scan_mode == "google_ocr":
        logger.info("PDF text extraction: Google OCR")
        return _extract_text_from_pdf_google_ocr(file_bytes)
    if scan_mode == "upstage_ocr":
        logger.info("PDF text extraction: Upstage OCR")
        return _extract_text_from_pdf_upstage_ocr(file_bytes)
    logger.info("PDF text extraction: PyMuPDF (text layer)")
    return _extract_text_from_pdf_pymupdf(file_bytes)


def _extract_text_from_pdf_pymupdf(file_bytes: bytes, max_pages: int | None = None) -> str:
    """Extract text from PDF using PyMuPDF. Pages separated by \\n\\n."""
    doc = fitz.open(stream=file_bytes, filetype="pdf")  # type: ignore[no-untyped-call]
    page_texts: list[str] = []

    total_pages: int = int(doc.page_count)  # type: ignore[arg-type]
    page_limit = min(total_pages, max_pages) if max_pages else total_pages

    for i in range(page_limit):
        page = doc[i]  # type: ignore[index]
        t = str(page.get_text()).strip()  # type: ignore[no-untyped-call]
        if t:
            page_texts.append(t)

    doc.close()  # type: ignore[no-untyped-call]
    return "\n\n".join(page_texts)


def _extract_text_from_pdf_upstage_ocr(file_bytes: bytes) -> str:
    """
    Extract text from PDF using Upstage Document OCR API.

    Endpoint: https://api.upstage.ai/v1/document-digitization
    Model: ocr
    """
    api_key = settings.UPSTAGE_API_KEY
    if not api_key:
        raise ValueError("UPSTAGE_API_KEY is not configured. Cannot perform OCR extraction.")

    url = "https://api.upstage.ai/v1/document-digitization"
    headers = {"Authorization": f"Bearer {api_key}"}

    response = httpx.post(
        url,
        headers=headers,
        files={"document": ("document.pdf", file_bytes, "application/pdf")},
        data={"model": "ocr"},
        timeout=300.0,
    )
    if response.status_code == 413:
        size_mb = len(file_bytes) / (1024 * 1024)
        raise ValueError(f"PDF file too large for OCR processing ({size_mb:.1f} MB). Please use a smaller file.")

    response.raise_for_status()

    result = response.json()

    if "text" in result:
        return result["text"]

    if "pages" in result:
        page_texts: list[str] = []
        for page in result["pages"]:
            text = page.get("text", "").strip()
            if text:
                page_texts.append(text)
        return "\n\n".join(page_texts)

    logger.warning(f"Unexpected OCR response structure: {list(result.keys())}")
    return str(result)


def _extract_text_from_pdf_google_ocr(file_bytes: bytes) -> str:
    """
    Extract text from PDF using Google Document AI (Enterprise Document OCR).

    Sends the entire PDF to the Document AI processor in a single request.
    Auth is handled via the GCP_DOCUMENTAI_KEY env var (JSON string) or
    falls back to GOOGLE_APPLICATION_CREDENTIALS (file path).

    API reference: https://docs.cloud.google.com/document-ai/docs/send-request
    """
    project_id = settings.GOOGLE_CLOUD_PROJECT_ID
    location = settings.GOOGLE_CLOUD_LOCATION
    processor_id = settings.GOOGLE_DOCUMENT_AI_PROCESSOR_ID

    if not project_id or not processor_id:
        raise ValueError(
            "GOOGLE_CLOUD_PROJECT_ID and GOOGLE_DOCUMENT_AI_PROCESSOR_ID must be configured for Google OCR extraction."
        )

    credentials = None
    if settings.GCP_DOCUMENTAI_KEY:
        credentials = service_account.Credentials.from_service_account_info(  # type: ignore[no-untyped-call]
            json.loads(settings.GCP_DOCUMENTAI_KEY)
        )

    opts = ClientOptions(api_endpoint=f"{location}-documentai.googleapis.com")
    client = documentai.DocumentProcessorServiceClient(client_options=opts, credentials=credentials)
    name = client.processor_path(project_id, location, processor_id)

    raw_document = documentai.RawDocument(
        content=file_bytes,
        mime_type="application/pdf",
    )

    request = documentai.ProcessRequest(
        name=name,
        raw_document=raw_document,
        field_mask="text",
    )

    result = client.process_document(request=request)  # type: ignore[no-untyped-call]
    return result.document.text


# =============================================================================
# RTF  (striprtf)
# =============================================================================


def extract_text_from_rtf(file_bytes: bytes) -> str:
    """Extract text from RTF. Paragraphs separated by \\n\\n."""

    rtf_content: str = decode_bytes(file_bytes)
    plain: str = cast(str, rtf_to_text(rtf_content, encoding="utf-8", errors="replace"))

    paragraphs: list[str] = [p.strip() for p in plain.split("\n") if p.strip()]
    return "\n\n".join(paragraphs)


# =============================================================================
# EPUB  (stub)
# =============================================================================


def extract_text_from_epub(file_bytes: bytes) -> str:
    """Extract text from EPUB. To be implemented."""
    raise NotImplementedError("EPUB text extraction is not yet implemented")


# =============================================================================
# Video  (stub)
# =============================================================================


def extract_text_from_video(file_url: str) -> str:
    """Extract subtitles/text from video. To be implemented."""
    raise NotImplementedError("Video text extraction is not yet implemented")
